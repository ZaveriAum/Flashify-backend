from openai import OpenAI
from config import Config
from decorators import AuthDecorator
from services.flashcard_service import FlashCardService
from services.note_service import NoteService
from utils.error import AppError
import json
import PyPDF2
from docx import Document
from pptx import Presentation
from io import BytesIO

client = OpenAI(api_key=Config.OPENAI_KEY)

class AIService:

    @staticmethod
    def extract_text(text=None, file=None):
        try:
            
            if file:
                    file_obj = file.get('file')
                    
                    if not file_obj:
                        raise AppError("No file provided.", 400)

                    file_stream = BytesIO(file_obj.read())
                    filename = file_obj.filename.lower()

                    if filename.endswith('.pdf'):
                        text = AIService.extract_from_pdf(file_stream)
                    elif filename.endswith('.docx'):
                        text = AIService.extract_from_docx(file_stream)
                    elif filename.endswith('.pptx'):
                        text = AIService.extract_from_pptx(file_stream)
                    elif filename.endswith('.txt'):
                        text = file_stream.getvalue().decode('utf-8')
                    else:
                        raise AppError("Unsupported file type", 400)
                
            if not text:
                raise AppError("Either text or file must be provided", 400)
            return text
        except Exception as e:
            raise AppError(getattr(e, 'message', 'Unknown Error'), getattr(e, 'statusCode', 400))

    @staticmethod
    @AuthDecorator.jwt_auth
    def generate_flashcards(text=None, file=None):
        try:
            text = AIService.extract_text(text, file)            
            
            base = """
                You are an advanced AI specializing in creating high-quality flashcards for graduate and post-graduate level study. Your task is to generate exactly 10 flashcards based on the given topic while strictly following these guidelines:

                Generate exactly 10 flashcards, no more and no less.
                Ensure all questions are challenging, in-depth, and appropriate for graduate-level study.
                Avoid basic definitions or questions that can be answered with a single word taken directly from the topic.
                Answers must be concise, containing no more than five words while maintaining accuracy. If necessary, answers can be short phrases, but they must remain brief.
                All answers must be unique with no repetitions across the flashcards.
                Stick strictly to the topic provided in the prompt, ensuring relevance and depth.
                Cover a diverse range of subtopics to ensure comprehensive understanding.
                Include questions that test conceptual understanding, theoretical knowledge, practical applications, and critical thinking.
                Avoid yes/no questions or any that can be answered with a simple affirmation or negation.
                Return the flashcards in the following structured JSON format:

                {
                    "flashcards": [
                        {
                        "question": "A detailed and thought-provoking question related to the topic",
                        "answer": "A concise and precise answer"
                        }
                    ]
                }

                For example, if the topic is "Fundamentals of Psychology," avoid simple questions such as "Which branch studies the human mind?" Instead, structure it as follows:

                Question: "What cognitive bias describes the tendency to search for or interpret information in a way that confirms one's preexisting beliefs?"
                Answer: "Confirmation bias"

                Ensure that all generated questions and answers meet these criteria before returning the final output.
            """

            prompt=f"{base}\n{text}"
            
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                { "role": "user", "content": prompt },
                ],
                temperature=0
            );
            
            flashcards_json = response.choices[0].message.content.strip("```json\n").strip("```")
            return json.loads(flashcards_json)
        except Exception as e:
            raise AppError(getattr(e, 'message', 'Unknown Error'), getattr(e, 'statusCode', 400))
    
    @staticmethod
    @AuthDecorator.jwt_auth
    def generate_note(text=None, file=None):
        try:
            text = AIService.extract_text(text, file)            
            base = """
                You are an advanced AI specializing in generating high-quality notes for graduate and post-graduate level study. Your task is to generate concise, structured, and well-organized notes based strictly on the given topic while adhering to the following constraints:

                The notes must be clear, precise, and well-structured, ensuring easy readability and comprehension.
                Stick strictly to the provided text and topic. Do not introduce any external information or make assumptions.
                The notes must not exceed 4000 characters in length.
                The output must only contain plain text—no images, emojis, formatting elements, or special characters.
                The notes should be informative and well-organized, covering key concepts, theories, applications, and critical points relevant to the topic.
                Ensure logical flow and coherence so that the notes provide a structured understanding of the subject.
                Avoid unnecessary repetition or overly simplistic points that do not add value.
                The output should not contain summaries of unrelated topics, opinions, or any content beyond the provided text.

                Your response should consist of only the generated notes, with no additional explanations or formatting outside the required text.
            """

            prompt=f"{base}\n{text}"
            
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                { "role": "user", "content": prompt },
                ],
                temperature=0
            );
            return response.choices[0].message.content.strip()
        except Exception as e:
            raise AppError(getattr(e, 'message', 'Unknown Error'), getattr(e, 'statusCode', 400))

    @staticmethod
    @AuthDecorator.jwt_auth
    def interact_with_flashcard(folder_id, prompt=None):
        try:
            flashcards = FlashCardService.get_flashcards(folder_id)
            flashcards = [(flashcard["question"], flashcard["answer"]) for flashcard in flashcards]
            if not prompt:
                raise AppError("Prompt empty", 400)
            else:
                prompt = f"""
                Role:
                    You are a distinguished professor at a prestigious Ivy League university. Your sole function is that of an academic chatbot—nothing more, nothing less.

                Context:
                    You will be provided with a set of flashcards and previous interactions with the user.

                    Flashcards:
                    {flashcards}

                    User’s Prompt:
                    {prompt}

                Task:
                    Your responsibility is to answer the question or complete the task specified in the prompt.

                Strict Instructions:
                    Do not include flashcard IDs, questions, or answers in your response.
                    Do not reference or repeat any flashcard content directly.
                    Provide a detailed, thorough, and well-explained answer.
                    Do not assume missing details; instead, focus on what is explicitly stated in the prompt.
                    Your response must be in clear, professional, and academic language that ensures full comprehension.
                    Output only text—no formatting, bullet points, or structural elements beyond what naturally fits in a well-written response.
                    Now, proceed with the task as described in the user’s prompt.
                """
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                { "role": "user", "content": prompt },
                ],
                temperature=0,
            );
            return response.choices[0].message.content
        except Exception as e:
            raise AppError(getattr(e, 'message', 'Unknown Error'), getattr(e, 'statusCode', 400))
    
    @staticmethod
    @AuthDecorator.jwt_auth
    def interact_with_note(note_id, prompt=None):
        try:
            note = NoteService.get_note_content(note_id)
            if not prompt:
                raise AppError("Prompt empty", 400)
            else:
                prompt = f"""
                Role:
                    You are a distinguished professor at a prestigious Ivy League university. Your sole function is that of an academic chatbot—nothing more, nothing less.

                Context:
                    You will be provided with a set of flashcards and previous interactions with the user.

                    Note Content:
                    {note['note']}

                    User’s Prompt:
                    {prompt}

                Task:
                    Your responsibility is to answer the question or complete the task specified in the prompt.

                Strict Instructions:
                    Do not repeat note content in your response.
                    Do not reference or repeat any note content directly.
                    Provide a detailed, thorough, and well-explained answer.
                    Do not assume missing details; instead, focus on what is explicitly stated in the prompt.
                    Your response must be in clear, professional, and academic language that ensures full comprehension.
                    Output only text—no formatting, bullet points, or structural elements beyond what naturally fits in a well-written response.
                    Now, proceed with the task as described in the user’s prompt.
                """
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                { "role": "user", "content": prompt },
                ],
                temperature=0,
            );
            return response.choices[0].message.content
        except Exception as e:
            raise AppError(getattr(e, 'message', 'Unknown Error'), getattr(e, 'statusCode', 400))
    
    # File extraction methods
    @staticmethod
    def extract_from_pdf(file_stream):
        reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    @staticmethod
    def extract_from_docx(file_stream):
        doc = Document(file_stream)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    @staticmethod
    def extract_from_pptx(file_stream):
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text